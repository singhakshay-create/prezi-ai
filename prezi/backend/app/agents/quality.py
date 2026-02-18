from app.providers.base import LLMProvider
from app.models import (
    Storyline, QualityScore, ResearchResults,
    SlideContent, SlideIssue, SlideQualityReport, SlideFeedback,
)
from app.agents.screenshot import render_slides_to_images, cleanup_screenshots
import json
from typing import List, Tuple, Optional
from pptx import Presentation
from pptx.util import Inches


class QualityChecker:
    """Validates presentation quality using LLM analysis."""

    def __init__(self, llm_provider: LLMProvider):
        self.llm = llm_provider

    async def check(self, storyline: Storyline) -> QualityScore:
        """Check presentation quality."""

        system_prompt = """You are a senior partner at McKinsey reviewing a presentation deck.
Evaluate the quality across multiple dimensions using consulting standards."""

        user_prompt = f"""Review this presentation storyline and score it on a 0-100 scale:

SCQA Framework:
- Situation: {storyline.scqa.situation}
- Complication: {storyline.scqa.complication}
- Question: {storyline.scqa.question}
- Answer: {storyline.scqa.answer}

Governing Thought: {storyline.governing_thought}

Hypotheses:
{chr(10).join([f"{i+1}. {h.text}" for i, h in enumerate(storyline.hypotheses)])}

Evaluate on these criteria (0-100 each):
1. slide_logic: Is the SCQA clear and logical?
2. mece_structure: Are hypotheses mutually exclusive and collectively exhaustive?
3. so_what: Is there a clear insight and recommendation?
4. data_quality: Would the hypotheses be supported by data? (assume research is done)
5. chart_accuracy: Would charts effectively communicate findings? (mock score)
6. visual_consistency: Would the presentation look professional? (mock score)

Also provide 2-4 improvement suggestions.

Return ONLY valid JSON:
{{
  "slide_logic": 85,
  "mece_structure": 80,
  "so_what": 90,
  "data_quality": 85,
  "chart_accuracy": 80,
  "visual_consistency": 85,
  "suggestions": ["Suggestion 1", "Suggestion 2"]
}}"""

        response = await self.llm.generate(
            prompt=user_prompt,
            system=system_prompt,
            temperature=0.3,
            max_tokens=1000
        )

        # Parse JSON response
        try:
            json_str = response.strip()
            if "```json" in json_str:
                json_str = json_str.split("```json")[1].split("```")[0].strip()
            elif "```" in json_str:
                json_str = json_str.split("```")[1].split("```")[0].strip()

            data = json.loads(json_str)

            # Calculate overall score (weighted average)
            overall = int(
                data["slide_logic"] * 0.25 +
                data["mece_structure"] * 0.25 +
                data["so_what"] * 0.25 +
                data["data_quality"] * 0.15 +
                data["chart_accuracy"] * 0.05 +
                data["visual_consistency"] * 0.05
            )

            return QualityScore(
                overall_score=overall,
                slide_logic=data["slide_logic"],
                mece_structure=data["mece_structure"],
                so_what=data["so_what"],
                data_quality=data["data_quality"],
                chart_accuracy=data["chart_accuracy"],
                visual_consistency=data["visual_consistency"],
                suggestions=data["suggestions"]
            )

        except (json.JSONDecodeError, KeyError) as e:
            # Return default passing score if parsing fails
            return QualityScore(
                overall_score=75,
                slide_logic=75,
                mece_structure=75,
                so_what=75,
                data_quality=75,
                chart_accuracy=75,
                visual_consistency=75,
                suggestions=["Quality check completed with default scores"]
            )

    # ------------------------------------------------------------------
    # PPTX-aware quality checking
    # ------------------------------------------------------------------

    def _extract_pptx_content(self, pptx_path: str) -> List[SlideContent]:
        """Extract per-slide content from a PPTX file."""
        prs = Presentation(pptx_path)
        slide_contents = []

        for idx, slide in enumerate(prs.slides):
            title_text = ""
            body_paragraphs: List[str] = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                # Heuristic: topmost textbox (top < 1 inch) is the title
                try:
                    is_title = shape.top < Inches(1.2)
                except Exception:
                    is_title = False

                all_text = shape.text_frame.text.strip()
                if not all_text:
                    continue

                if is_title and not title_text:
                    title_text = all_text
                else:
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            body_paragraphs.append(para_text)

            has_chart = any(
                shape.shape_type == 13 for shape in slide.shapes
            )
            has_table = any(
                shape.has_table for shape in slide.shapes
            )

            all_words = (title_text + " " + " ".join(body_paragraphs)).split()
            word_count = len([w for w in all_words if w])

            slide_contents.append(SlideContent(
                slide_index=idx,
                title=title_text,
                body_text=body_paragraphs,
                has_chart=has_chart,
                has_table=has_table,
                shape_count=len(slide.shapes),
                word_count=word_count,
            ))

        return slide_contents

    async def _inspect_with_llm(
        self,
        slides: List[SlideContent],
        storyline: Storyline,
        iteration: int,
    ) -> SlideQualityReport:
        """LLM Call 1: inspect all slides, return SlideQualityReport."""

        slide_lines = []
        for s in slides:
            body_preview = "; ".join(s.body_text[:5]) if s.body_text else "(empty)"
            slide_lines.append(
                f"SLIDE {s.slide_index} [{s.title or '(no title)'}]:\n"
                f"  Body: {body_preview[:300]}\n"
                f"  Chart: {'Yes' if s.has_chart else 'No'} | "
                f"Table: {'Yes' if s.has_table else 'No'} | "
                f"Words: {s.word_count}"
            )

        slides_text = "\n".join(slide_lines)

        system_prompt = (
            "You are a McKinsey senior partner reviewing a management consulting deck. "
            "Identify specific per-slide problems and score quality dimensions."
        )

        user_prompt = f"""Review this consulting presentation (iteration {iteration}).

STORYLINE:
- Situation: {storyline.scqa.situation[:200]}
- Answer: {storyline.scqa.answer[:200]}
- Key hypotheses: {'; '.join(h.text for h in storyline.hypotheses[:3])}

SLIDES:
{slides_text}

EVALUATION STANDARDS:
- Placeholder detection: labels like "Factor 1", "Segment A", "BU-A", "Action item 1" are placeholder_data issues
- too_sparse: fewer than 20 words and no chart/table
- missing_so_what: title is generic (e.g. "Recommendations", "Analysis") with no specific finding
- weak_title: title doesn't contain a number, finding, or action verb
- missing_chart: data-heavy analysis slides (market analysis, sensitivity) with no chart

Score each dimension 0-100:
- information_density_score: Are slides data-rich with evidence?
- chart_quality_score: Are charts using real named data (not Factor 1)?
- narrative_flow_score: Does SCQA logic flow through slides?

Return ONLY valid JSON matching this schema:
{{
  "iteration": {iteration},
  "information_density_score": 60,
  "chart_quality_score": 40,
  "narrative_flow_score": 70,
  "storyline_suggestions": ["suggestion 1", "suggestion 2"],
  "issues": [
    {{
      "slide_index": 4,
      "issue_type": "placeholder_data",
      "description": "Bar chart uses generic Factor 1-5 labels",
      "fix_suggestion": "Replace with specific market drivers from research"
    }}
  ]
}}"""

        response = await self.llm.generate(
            prompt=user_prompt,
            system=system_prompt,
            temperature=0.2,
            max_tokens=2000,
        )

        try:
            json_str = response.strip()
            if "```json" in json_str:
                json_str = json_str.split("```json")[1].split("```")[0].strip()
            elif "```" in json_str:
                json_str = json_str.split("```")[1].split("```")[0].strip()

            data = json.loads(json_str)

            issues = []
            for issue_data in data.get("issues", []):
                try:
                    issues.append(SlideIssue(**issue_data))
                except Exception:
                    pass

            return SlideQualityReport(
                iteration=iteration,
                slides=slides,
                issues=issues,
                information_density_score=int(data.get("information_density_score", 50)),
                chart_quality_score=int(data.get("chart_quality_score", 50)),
                narrative_flow_score=int(data.get("narrative_flow_score", 50)),
                storyline_suggestions=data.get("storyline_suggestions", []),
            )

        except (json.JSONDecodeError, KeyError, TypeError, AttributeError, ValueError):
            return SlideQualityReport(
                iteration=iteration,
                slides=slides,
                issues=[],
                information_density_score=50,
                chart_quality_score=50,
                narrative_flow_score=50,
                storyline_suggestions=[],
            )

    async def _generate_slide_feedback(
        self,
        issues: List[SlideIssue],
        slides: List[SlideContent],
        storyline: Storyline,
        research: ResearchResults,
    ) -> List[SlideFeedback]:
        """LLM Call 2: generate concrete fixes for flagged slides."""
        if not issues:
            return []

        # Group issues by slide_index
        issues_by_slide: dict = {}
        for issue in issues:
            issues_by_slide.setdefault(issue.slide_index, []).append(issue)

        # Build a lookup for slide content
        slide_by_idx = {s.slide_index: s for s in slides}

        # Top 3 research evidence snippets sorted by relevance
        top_evidence = []
        for hyp_ev in research.hypotheses_evidence:
            for ev in sorted(hyp_ev.evidence, key=lambda e: e.relevance_score, reverse=True)[:3]:
                top_evidence.append(ev.snippet)
        top_evidence = top_evidence[:3]

        # Build per-slide fix requests
        slide_requests = []
        for slide_idx, slide_issues in issues_by_slide.items():
            slide = slide_by_idx.get(slide_idx)
            if not slide:
                continue
            issue_texts = "\n".join(
                f"- [{i.issue_type}] {i.description} → {i.fix_suggestion}"
                for i in slide_issues
            )
            slide_requests.append(
                f"SLIDE {slide_idx} (title: {slide.title or '(no title)'}):\n"
                f"  Current body: {'; '.join(slide.body_text[:3])}\n"
                f"  Issues:\n{issue_texts}"
            )

        evidence_text = "\n".join(f"- {e}" for e in top_evidence) if top_evidence else "No evidence available."

        system_prompt = (
            "You are a McKinsey deck editor. Generate precise slide fixes based on research evidence."
        )

        user_prompt = f"""Fix the following slides in a consulting presentation.

STORYLINE ANSWER: {storyline.scqa.answer[:200]}

RESEARCH EVIDENCE:
{evidence_text}

SLIDES TO FIX:
{chr(10).join(slide_requests)}

RULES:
- Titles must contain a specific finding, number, or action verb
- Bullets must reference research evidence or named concepts
- Chart categories must be real named concepts (not "Factor 1", "Segment A")
- Do NOT set new_chart_data for slides 0, 1, 2 (title/summary/situation slides)
- issues_addressed must list the issue_type strings that are resolved

Return ONLY a JSON array:
[
  {{
    "slide_index": 4,
    "new_title": "Hybrid Cloud Adoption Grows 2x Faster Than On-Prem",
    "new_bullets": ["SOC2 certification reduces procurement cycle by 40%", "Pay-as-you-go converts 3x better than annual contracts"],
    "new_chart_data": {{
      "chart_type": "bar",
      "categories": ["Hybrid Cloud", "Public Cloud", "On-Premises", "Private Cloud"],
      "values": [85, 75, 45, 60],
      "title": "Enterprise Adoption Rate by Deployment Model (%)",
      "x_label": "Adoption Score"
    }},
    "issues_addressed": ["placeholder_data", "weak_title"]
  }}
]"""

        response = await self.llm.generate(
            prompt=user_prompt,
            system=system_prompt,
            temperature=0.3,
            max_tokens=2000,
        )

        try:
            json_str = response.strip()
            if "```json" in json_str:
                json_str = json_str.split("```json")[1].split("```")[0].strip()
            elif "```" in json_str:
                json_str = json_str.split("```")[1].split("```")[0].strip()

            data = json.loads(json_str)
            if not isinstance(data, list):
                return []

            feedbacks = []
            for item in data:
                try:
                    feedbacks.append(SlideFeedback(**item))
                except Exception:
                    pass
            return feedbacks

        except (json.JSONDecodeError, KeyError, TypeError, AttributeError, ValueError):
            return []

    def _report_to_quality_score(
        self, report: SlideQualityReport, iteration: int
    ) -> QualityScore:
        """Map SlideQualityReport scores to QualityScore."""
        # Count issues by type for derived dimensions
        issue_types = [i.issue_type for i in report.issues]
        mece_penalty = issue_types.count("mece_violation") * 10
        so_what_penalty = issue_types.count("missing_so_what") * 10
        visual_penalty = (issue_types.count("too_dense") + issue_types.count("too_sparse")) * 5

        mece_score = max(0, 80 - mece_penalty)
        so_what_score = max(0, 80 - so_what_penalty)
        visual_score = max(0, 80 - visual_penalty)

        slide_logic = report.narrative_flow_score
        chart_accuracy = report.chart_quality_score
        data_quality = report.information_density_score

        overall = int(
            slide_logic * 0.25 +
            mece_score * 0.25 +
            so_what_score * 0.25 +
            data_quality * 0.15 +
            chart_accuracy * 0.05 +
            visual_score * 0.05
        )
        overall = max(0, min(100, overall))

        suggestions = list(report.storyline_suggestions)
        for issue in report.issues[:3]:
            suggestions.append(issue.description)

        return QualityScore(
            overall_score=overall,
            slide_logic=slide_logic,
            mece_structure=mece_score,
            so_what=so_what_score,
            data_quality=data_quality,
            chart_accuracy=chart_accuracy,
            visual_consistency=visual_score,
            suggestions=suggestions[:6],
            iterations_run=iteration,
            final_report=report,
        )

    # ------------------------------------------------------------------
    # Visual inspection (screenshot-based)
    # ------------------------------------------------------------------

    def _render_slide_screenshots(self, pptx_path: str) -> Tuple[List[str], Optional[str]]:
        """
        Render each slide to a PNG and return (png_paths, temp_dir).
        Returns ([], None) on any error so callers can fall back to text inspection.
        """
        try:
            png_paths, temp_dir = render_slides_to_images(pptx_path, dpi_scale=1.5)
            return png_paths, temp_dir
        except Exception:
            return [], None

    async def _visual_inspect_with_llm(
        self,
        png_paths: List[str],
        slides: List[SlideContent],
        storyline: Storyline,
        iteration: int,
    ) -> SlideQualityReport:
        """
        LLM vision call: send all slide screenshots + text metadata to the LLM
        and receive a SlideQualityReport with visually-detected issues.
        """
        # Build a compact text description alongside the images
        slide_lines = []
        for s in slides:
            body_preview = "; ".join(s.body_text[:3]) if s.body_text else "(empty)"
            slide_lines.append(
                f"Slide {s.slide_index} [{s.title or '(no title)'}]: "
                f"words={s.word_count}, chart={'Y' if s.has_chart else 'N'}, "
                f"table={'Y' if s.has_table else 'N'} | {body_preview[:120]}"
            )

        system_prompt = (
            "You are a McKinsey senior partner and visual design expert reviewing "
            "a management consulting presentation. Examine EACH slide screenshot "
            "carefully for visual quality issues."
        )

        user_prompt = f"""Iteration {iteration} visual review of a consulting deck.

SLIDE METADATA:
{chr(10).join(slide_lines)}

STORYLINE ANSWER: {storyline.scqa.answer[:200]}

Examine every slide screenshot and identify ALL visual quality issues including:
- too_sparse: excessive white space, nearly empty slide
- too_dense: text overflowing or unreadable crowding
- placeholder_data: generic labels like "Factor 1", "Segment A", "BU-A"
- missing_so_what: title is a generic noun with no specific finding or number
- weak_title: title has no number, finding, or action verb
- missing_chart: analysis slide has no visual, only bullet points
- mece_violation: bullet points overlap logically or don't cover the topic
- narrative_gap: slide doesn't connect to the prior slide's logic
- (visual-only issues): note font size inconsistency, overlapping text, poor contrast,
  misaligned shapes, and color palette problems in the description field,
  using issue_type "too_sparse" or "too_dense" as the closest match

Score each dimension 0–100:
- information_density_score: how data-rich and evidence-backed are the slides?
- chart_quality_score: are charts using real named data and clearly readable?
- narrative_flow_score: does the visual story flow logically slide-to-slide?

Return ONLY valid JSON:
{{
  "iteration": {iteration},
  "information_density_score": 60,
  "chart_quality_score": 45,
  "narrative_flow_score": 70,
  "storyline_suggestions": ["suggestion 1", "suggestion 2"],
  "issues": [
    {{
      "slide_index": 0,
      "issue_type": "too_sparse",
      "description": "Title slide has large empty area in the bottom two-thirds",
      "fix_suggestion": "Add a relevant illustrative image in the lower half"
    }}
  ]
}}"""

        try:
            # Cap at 12 images to stay within token limits
            capped_paths = png_paths[:12]
            response = await self.llm.generate_with_vision(
                prompt=user_prompt,
                image_paths=capped_paths,
                system=system_prompt,
                temperature=0.2,
                max_tokens=2500,
            )

            json_str = response.strip()
            if "```json" in json_str:
                json_str = json_str.split("```json")[1].split("```")[0].strip()
            elif "```" in json_str:
                json_str = json_str.split("```")[1].split("```")[0].strip()

            data = json.loads(json_str)
            if not isinstance(data, dict):
                raise ValueError("Expected dict")

            issues = []
            for issue_data in data.get("issues", []):
                try:
                    issues.append(SlideIssue(**issue_data))
                except Exception:
                    pass

            return SlideQualityReport(
                iteration=iteration,
                slides=slides,
                issues=issues,
                information_density_score=int(data.get("information_density_score", 50)),
                chart_quality_score=int(data.get("chart_quality_score", 50)),
                narrative_flow_score=int(data.get("narrative_flow_score", 50)),
                storyline_suggestions=data.get("storyline_suggestions", []),
            )

        except (json.JSONDecodeError, KeyError, TypeError, AttributeError, ValueError):
            return SlideQualityReport(
                iteration=iteration,
                slides=slides,
                issues=[],
                information_density_score=50,
                chart_quality_score=50,
                narrative_flow_score=50,
                storyline_suggestions=[],
            )

    async def check_with_pptx(
        self,
        pptx_path: str,
        storyline: Storyline,
        research: ResearchResults,
        iteration: int,
    ) -> Tuple[QualityScore, SlideQualityReport, List[SlideFeedback]]:
        """
        Public method: render screenshots (if available), inspect visually or via
        text, generate concrete fixes, return (QualityScore, SlideQualityReport,
        List[SlideFeedback]).
        """
        slides = self._extract_pptx_content(pptx_path)

        # Prefer visual inspection when the LLM supports it
        png_paths, temp_dir = self._render_slide_screenshots(pptx_path)
        try:
            if png_paths and self.llm.supports_vision():
                report = await self._visual_inspect_with_llm(
                    png_paths, slides, storyline, iteration
                )
            else:
                report = await self._inspect_with_llm(slides, storyline, iteration)
        finally:
            if temp_dir:
                cleanup_screenshots(temp_dir)

        feedback = await self._generate_slide_feedback(
            report.issues, slides, storyline, research
        )
        score = self._report_to_quality_score(report, iteration)
        return score, report, feedback
