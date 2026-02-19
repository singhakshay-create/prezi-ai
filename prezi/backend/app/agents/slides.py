from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from app.models import Storyline, ResearchResults, Hypothesis, HypothesisEvidence
from typing import Literal, Optional, List, Tuple
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
import io
import os
import re
from datetime import datetime


def _coerce_float(v) -> float:
    """Return a float from v, extracting from a dict if necessary."""
    if isinstance(v, (int, float)):
        return float(v)
    if isinstance(v, dict):
        for val in v.values():
            try:
                return float(val)
            except (TypeError, ValueError):
                pass
        return 0.0
    try:
        return float(v)
    except (TypeError, ValueError):
        return 0.0


def _coerce_str(v) -> str:
    """Return a plain string from v, extracting a label field from a dict if present."""
    if isinstance(v, str):
        return v
    if isinstance(v, dict):
        for key in ("label", "name", "text", "category", "title"):
            if key in v:
                return str(v[key])
        # Fall back to first string value found
        for val in v.values():
            if isinstance(val, str):
                return val
        return str(v)
    return str(v)


class SlideGenerator:
    """Generates consulting-style presentations using python-pptx."""

    def __init__(self, template_path: str = None):
        self.template_path = template_path
        # McKinsey colors
        self.primary_color = RGBColor(0, 51, 153)  # McKinsey blue
        self.accent_color = RGBColor(0, 176, 240)  # Light blue
        self._last_pptx_path: str = None

    async def create_presentation(
        self,
        topic: str,
        storyline: Storyline,
        research: ResearchResults,
        length: Literal["short", "medium", "long"]
    ) -> str:
        """Create PPTX presentation."""

        # Create presentation (use template if provided)
        if self.template_path and os.path.isfile(self.template_path):
            prs = Presentation(self.template_path)
        else:
            prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Short base: title + situation + complication + one slide per hypothesis
        self._add_title_slide(prs, topic, storyline)
        self._add_situation_slide(prs, storyline)
        self._add_complication_slide(prs, storyline)
        for hyp in storyline.hypotheses:
            evidence = next(
                (e for e in research.hypotheses_evidence if e.hypothesis_id == hyp.id),
                None
            )
            self._add_hypothesis_slide(prs, hyp, evidence)

        # Medium: add 4 standalone chart slides
        if length in ["medium", "long"]:
            self._add_bar_chart_slide(prs, storyline, research)
            self._add_waterfall_slide(prs, storyline)
            self._add_pie_chart_slide(prs, storyline, research)
            self._add_tornado_chart_slide(prs, storyline, research)

        # Long: add 5 framework slides
        if length == "long":
            self._add_marimekko_chart_slide(prs, storyline, research)
            self._add_bcg_matrix_slide(prs, storyline, research)
            self._add_priority_matrix_slide(prs, storyline)
            self._add_value_chain_slide(prs, storyline)
            self._add_heatmap_slide(prs, storyline, research)

        self._add_recommendations(prs, storyline)
        self._add_sources(prs, research)

        # Save presentation
        os.makedirs("./data/presentations", exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"presentation_{timestamp}.pptx"
        filepath = f"./data/presentations/{filename}"

        prs.save(filepath)
        self._last_pptx_path = filepath
        return filepath

    # ------------------------------------------------------------------
    # Chrome helper — used by ALL content slides
    # ------------------------------------------------------------------

    def _add_slide_title(self, slide, title: str):
        """Add consulting chrome (left stripe, title, separator) to a slide."""
        # Thick dark blue left accent stripe (0.25" for professional look)
        stripe = slide.shapes.add_shape(1, 0, 0, Inches(0.25), Inches(7.5))
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = RGBColor(0, 51, 153)
        stripe.line.fill.background()

        # Title textbox — wider for 13.333" widescreen
        title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.7), Inches(0.7))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = self._strip_markdown(title)
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 153)

        # Accent blue separator line (not gray)
        sep = slide.shapes.add_shape(1, Inches(0.4), Inches(0.88), Inches(12.6), Inches(0.02))
        sep.fill.solid()
        sep.fill.fore_color.rgb = RGBColor(0, 176, 240)
        sep.line.fill.background()

    # ------------------------------------------------------------------
    # Footer helper — used by all content slides
    # ------------------------------------------------------------------

    def _add_footer(self, slide, source: str = None):
        """Thin gray rule + 8pt source citation at bottom of slide."""
        rule = slide.shapes.add_shape(1, Inches(0.3), Inches(7.1), Inches(12.7), Inches(0.015))
        rule.fill.solid()
        rule.fill.fore_color.rgb = RGBColor(204, 204, 204)
        rule.line.fill.background()

        src_box = slide.shapes.add_textbox(Inches(0.3), Inches(7.15), Inches(12.7), Inches(0.3))
        tf = src_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Source: {source}" if source else "Source: Market research and industry data"
        p.font.size = Pt(8)
        p.font.italic = True
        p.font.color.rgb = RGBColor(128, 128, 128)

    # ------------------------------------------------------------------
    # KEY INSIGHT sidebar helper
    # ------------------------------------------------------------------

    def _add_insight_sidebar(self, slide, headline: str, bullets: List[str],
                              top_label: str = None, top_value: str = None):
        """Right sidebar: accent-bordered callout with headline, large metric, and bullets.

        Layout (13.333" widescreen slide):
          chart area    = 0.4" → 9.0"  (width 8.6")
          gap           = 0.2"
          sidebar left  = 9.2"
          sidebar width = 3.8"
          sidebar top   = 1.15"
          sidebar height = 5.75"  (stops above footer rule at 7.1")
        """
        SIDEBAR_L = Inches(9.2)
        SIDEBAR_W = Inches(3.8)
        SIDEBAR_T = Inches(1.15)
        SIDEBAR_H = Inches(5.75)
        PAD = Inches(0.15)   # inner padding from box edge

        # Cream container with accent border
        box = slide.shapes.add_shape(1, SIDEBAR_L, SIDEBAR_T, SIDEBAR_W, SIDEBAR_H)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 250, 240)   # cream
        box.line.color.rgb = RGBColor(0, 176, 240)           # accent blue border
        box.line.width = Pt(1.5)

        # "KEY INSIGHT" label
        hdr = slide.shapes.add_textbox(
            SIDEBAR_L + PAD, SIDEBAR_T + PAD, SIDEBAR_W - 2 * PAD, Inches(0.35))
        p = hdr.text_frame.paragraphs[0]
        p.text = "KEY INSIGHT"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 176, 240)

        # Optional large metric (top_label + top_value)
        y_offset = SIDEBAR_T + PAD + Inches(0.35)
        if top_label and top_value:
            metric_box = slide.shapes.add_textbox(
                SIDEBAR_L + PAD, y_offset, SIDEBAR_W - 2 * PAD, Inches(1.1))
            tf = metric_box.text_frame
            tf.word_wrap = True
            p_lbl = tf.paragraphs[0]
            p_lbl.text = top_label
            p_lbl.font.size = Pt(9)
            p_lbl.font.color.rgb = RGBColor(80, 80, 80)
            p_val = tf.add_paragraph()
            p_val.text = top_value
            p_val.font.size = Pt(28)
            p_val.font.bold = True
            p_val.font.color.rgb = RGBColor(0, 51, 153)
            y_offset += Inches(1.1)

        # Thin rule between metric and bullets
        if top_label and top_value:
            rule = slide.shapes.add_shape(
                1, SIDEBAR_L + PAD, y_offset, SIDEBAR_W - 2 * PAD, Inches(0.012))
            rule.fill.solid()
            rule.fill.fore_color.rgb = RGBColor(0, 176, 240)
            rule.line.fill.background()
            y_offset += Inches(0.08)

        # Insight bullets
        remaining_h = (SIDEBAR_T + SIDEBAR_H) - y_offset - PAD
        bul_box = slide.shapes.add_textbox(
            SIDEBAR_L + PAD, y_offset, SIDEBAR_W - 2 * PAD, remaining_h)
        tf = bul_box.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u25b6 {b}"     # ▶ triangle bullet
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(30, 30, 30)
            if i > 0:
                p.space_before = Pt(8)

    # ------------------------------------------------------------------
    # Sidebar content deriver
    # ------------------------------------------------------------------

    def _derive_sidebar_content(self, chart_data: dict, storyline_title: str):
        """Derive KEY INSIGHT sidebar content automatically from chart data."""
        categories = [_coerce_str(c) for c in chart_data.get("categories", [])]
        values     = [_coerce_float(v) for v in chart_data.get("values", [])]
        metric     = chart_data.get("x_label", "")
        if categories and values:
            idx       = values.index(max(values))
            top_label = f"Top driver ({metric})"
            top_value = f"{values[idx]:.0f}"
            bullets   = [
                f"{categories[idx]} leads at {values[idx]:.0f} {metric}",
                f"Average: {sum(values)/len(values):.0f} {metric}",
            ]
        else:
            top_label, top_value = None, None
            bullets = [storyline_title]
        return top_label, top_value, bullets

    # ------------------------------------------------------------------
    # Title slide
    # ------------------------------------------------------------------

    def _add_title_slide(self, prs, topic: str, storyline: Storyline):
        """Full-width dark navy title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Full-slide dark background
        bg = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0, 31, 96)
        bg.line.fill.background()

        # Thin accent line — full widescreen width
        accent = slide.shapes.add_shape(1, 0, Inches(2.3), Inches(13.333), Inches(0.04))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(0, 102, 255)
        accent.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.4))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = topic
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Governing thought
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(1.2))
        tf2 = sub_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = storyline.governing_thought
        p2.font.size = Pt(15)
        p2.font.color.rgb = RGBColor(153, 187, 255)

        # Date
        date_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(6.0), Inches(0.5))
        p3 = date_box.text_frame.paragraphs[0]
        p3.text = datetime.now().strftime("%B %Y")
        p3.font.size = Pt(11)
        p3.font.color.rgb = RGBColor(119, 153, 204)

    # ------------------------------------------------------------------
    # Situation slide
    # ------------------------------------------------------------------

    def _add_situation_slide(self, prs, storyline: Storyline):
        """Full-width situation slide with action title and finding bullets."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, storyline.scqa.situation_title or "Market Context")

        content_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.1), Inches(12.6), Inches(5.9))
        tf = content_box.text_frame
        tf.word_wrap = True

        bullets = (
            storyline.scqa.situation_bullets
            if storyline.scqa.situation_bullets
            else [s.strip() for s in storyline.scqa.situation.split(". ") if s.strip()]
        )

        self._render_finding_bullets(tf, bullets)
        self._add_footer(slide)

    # ------------------------------------------------------------------
    # Complication slide
    # ------------------------------------------------------------------

    def _add_complication_slide(self, prs, storyline: Storyline):
        """Full-width complication slide with action title and finding bullets."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, storyline.scqa.complication_title or "Key Challenges")

        content_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.1), Inches(12.6), Inches(5.9))
        tf = content_box.text_frame
        tf.word_wrap = True

        bullets = (
            storyline.scqa.complication_bullets
            if storyline.scqa.complication_bullets
            else [s.strip() for s in storyline.scqa.complication.split(". ") if s.strip()]
        )

        self._render_finding_bullets(tf, bullets)
        self._add_footer(slide)

    @staticmethod
    def _render_bold_text(paragraph, text: str):
        """Parse **term** markdown into PPTX runs with selective bold formatting."""
        parts = re.split(r'(\*\*[^*]+\*\*)', text)
        for part in parts:
            if part.startswith('**') and part.endswith('**'):
                run = paragraph.add_run()
                run.text = part[2:-2]
                run.font.bold = True
            elif part:
                run = paragraph.add_run()
                run.text = part
                run.font.bold = False

    @staticmethod
    def _strip_markdown(text: str) -> str:
        """Remove markdown formatting markers from text destined for plain PPTX paragraphs."""
        # Strip **bold** and *italic* markers
        text = re.sub(r'\*{1,2}([^*]+)\*{1,2}', r'\1', text)
        # Strip [link text](url) → link text
        text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
        return text.strip()

    def _render_finding_bullets(self, tf, bullets: List[str]):
        """Render bullets into a text frame: 12pt finding with bold **markers** + optional 9pt gray source."""
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            # Split on " — " to separate finding from source citation
            if " — " in bullet:
                finding, source_text = bullet.rsplit(" — ", 1)
            else:
                finding = bullet
                source_text = None

            self._render_bold_text(p, finding)
            p.font.size = Pt(12)
            if i > 0:
                p.space_before = Pt(10)

            if source_text:
                p_src = tf.add_paragraph()
                p_src.text = source_text
                p_src.font.size = Pt(9)
                p_src.font.color.rgb = RGBColor(128, 128, 128)

    def _add_native_bar_chart(self, slide, chart_data: dict, left, top, width, height):
        """Add a native editable PowerPoint BAR_CLUSTERED chart (users can edit data in Excel)."""
        categories = chart_data.get("categories", [f"Factor {i+1}" for i in range(5)])
        values = chart_data.get("values", [75, 85, 65, 90, 70])
        n = min(len(categories), len(values))
        categories = [_coerce_str(c) for c in categories[:n]]
        values = [_coerce_float(v) for v in values[:n]]

        cd = ChartData()
        cd.categories = categories
        cd.add_series(chart_data.get("x_label", "Values"), values)
        gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, cd)
        chart = gf.chart
        chart.has_legend = False
        series = chart.plots[0].series[0]
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = RGBColor(0, 51, 153)
        series.data_labels.show_value = True

    # ------------------------------------------------------------------
    # Hypothesis slide: chart left + evidence column right
    # ------------------------------------------------------------------

    def _add_hypothesis_slide(self, prs, hypothesis: Hypothesis, evidence):
        """Add a hypothesis slide: action title, bar chart left, light-blue evidence column right."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, hypothesis.action_title or hypothesis.text)

        # Build chart data from hint or use fallback
        if hypothesis.chart_hint:
            cats = hypothesis.chart_hint.get("categories", ["Baseline", "With Solution", "Industry Best"])
            n = len(cats)
            # Use LLM-provided values if available; otherwise fall back to formula
            if "values" in hypothesis.chart_hint:
                vals = [_coerce_float(v) for v in hypothesis.chart_hint["values"]]
            else:
                vals = [max(15, 85 - i * 20) for i in range(n)]
            chart_data = {
                "categories": cats,
                "values": vals,
                "title": hypothesis.text[:50],
                "x_label": hypothesis.chart_hint.get("metric", "Score"),
            }
        else:
            chart_data = {
                "categories": ["Baseline", "With Solution", "Industry Best"],
                "values": [45, 85, 70],
                "title": hypothesis.text[:50],
                "x_label": "Score",
            }

        self._add_native_bar_chart(slide, chart_data, Inches(0.4), Inches(1.15), Inches(7.0), Inches(5.75))

        # Pyramid level label — top-right corner
        level_box = slide.shapes.add_textbox(Inches(11.8), Inches(0.2), Inches(1.3), Inches(0.5))
        level_tf = level_box.text_frame
        lp = level_tf.paragraphs[0]
        lp.text = "HYPOTHESIS"
        lp.font.size = Pt(8)
        lp.font.bold = True
        lp.font.color.rgb = RGBColor(0, 176, 240)
        lp.alignment = PP_ALIGN.RIGHT

        # Light blue evidence column background (drawn before textbox for proper z-order)
        ev_bg = slide.shapes.add_shape(1, Inches(7.6), Inches(1.15), Inches(5.4), Inches(5.75))
        ev_bg.fill.solid()
        ev_bg.fill.fore_color.rgb = RGBColor(240, 248, 255)   # alice blue
        ev_bg.line.color.rgb = RGBColor(0, 176, 240)
        ev_bg.line.width = Pt(0.75)

        # Right column: evidence bullets
        bullets_box = slide.shapes.add_textbox(Inches(7.75), Inches(1.25), Inches(5.1), Inches(5.5))
        tf = bullets_box.text_frame
        tf.word_wrap = True

        ev_bullets = self._get_evidence_bullets(evidence)
        if ev_bullets:
            for i, (snippet, source) in enumerate(ev_bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = self._strip_markdown(snippet)
                p.font.size = Pt(11)
                p.font.bold = True
                if i > 0:
                    p.space_before = Pt(12)

                p_src = tf.add_paragraph()
                p_src.text = source
                p_src.font.size = Pt(9)
                p_src.font.color.rgb = RGBColor(128, 128, 128)

            # Confidence badge
            if evidence and evidence.confidence:
                p_conf = tf.add_paragraph()
                conf_color = {
                    "high": RGBColor(0, 128, 0),
                    "medium": RGBColor(255, 140, 0),
                    "low": RGBColor(200, 0, 0),
                }
                p_conf.text = f"Confidence: {evidence.confidence.capitalize()}"
                p_conf.font.size = Pt(9)
                p_conf.font.bold = True
                p_conf.space_before = Pt(14)
                p_conf.font.color.rgb = conf_color.get(evidence.confidence, RGBColor(128, 128, 128))
        else:
            p = tf.paragraphs[0]
            p.text = "Research evidence to be populated"
            p.font.size = Pt(11)

        # Footer with evidence source if available
        footer_source = None
        if evidence and evidence.evidence:
            footer_source = evidence.evidence[0].source
        self._add_footer(slide, footer_source)

    def _get_evidence_bullets(self, evidence, n: int = 3) -> List[Tuple[str, str]]:
        """Return top-N evidence (snippet, source) pairs sorted by relevance."""
        if not evidence or not evidence.evidence:
            return []
        sorted_ev = sorted(evidence.evidence, key=lambda e: e.relevance_score, reverse=True)
        return [(e.snippet[:120], e.source[:50]) for e in sorted_ev[:n]]

    # ------------------------------------------------------------------
    # Data-driven chart renderers
    # ------------------------------------------------------------------

    def _render_bar_chart(self, chart_data: dict) -> io.BytesIO:
        """Render a horizontal bar chart from a data dict and return BytesIO PNG."""
        categories = chart_data.get("categories", [f"Factor {i+1}" for i in range(5)])
        values = chart_data.get("values", [75, 85, 65, 90, 70])
        title = chart_data.get("title", "Key Success Factors")
        x_label = chart_data.get("x_label", "Impact Score")

        # Guard against LLM returning mismatched list lengths
        n = min(len(categories), len(values))
        categories, values = categories[:n], values[:n]

        # Coerce to plain types so matplotlib can hash/plot them
        categories = [_coerce_str(c) for c in categories]
        values = [_coerce_float(v) for v in values]

        fig, ax = plt.subplots(figsize=(8, 4))
        ax.barh(categories, values, color='#003399')
        ax.set_xlabel(x_label, fontsize=12)
        ax.set_title(title, fontsize=14, fontweight='bold')
        ax.grid(axis='x', alpha=0.3)

        # Add value labels on each bar
        max_val = max(values) if values else 1
        for i, val in enumerate(values):
            ax.text(val + max_val * 0.01, i, f"{val:.0f}",
                    va='center', fontsize=10, fontweight='bold', color='#003399')
        ax.set_xlim(0, max_val * 1.18)  # room for labels

        img_bytes = io.BytesIO()
        plt.savefig(img_bytes, format='png', bbox_inches='tight', dpi=150)
        img_bytes.seek(0)
        plt.close()
        return img_bytes

    def _render_waterfall_chart(self, chart_data: dict) -> io.BytesIO:
        """Render a waterfall chart from a data dict and return BytesIO PNG."""
        categories = chart_data.get("categories", ['Starting', 'Revenue\nGrowth', 'Cost\nReduction', 'Efficiency', 'Ending'])
        values = chart_data.get("values", [100, 25, 15, 10, 150])
        title = chart_data.get("title", "Value Creation Opportunity")

        # Guard against mismatched list lengths; waterfall needs at least 2 entries
        n = min(len(categories), len(values))
        n = max(n, 2)
        categories, values = categories[:n], values[:n]

        # Coerce to plain types so matplotlib can hash/plot them
        categories = [_coerce_str(c) for c in categories]
        values = [_coerce_float(v) for v in values]

        cumulative = []
        running = 0
        for i, v in enumerate(values):
            if i == 0 or i == len(values) - 1:
                running = v
            else:
                running += v
            cumulative.append(running)

        colors = ['#0033cc'] + ['#00b0f0'] * (len(values) - 2) + ['#0033cc']

        fig, ax = plt.subplots(figsize=(8, 4))
        for i, (cat, val, cum) in enumerate(zip(categories, values, cumulative)):
            if i == 0:
                ax.bar(i, val, color=colors[i], edgecolor='black', linewidth=1)
            elif i == len(categories) - 1:
                ax.bar(i, val, color=colors[i], edgecolor='black', linewidth=1)
            else:
                ax.bar(i, val, bottom=cumulative[i - 1], color=colors[i], edgecolor='black', linewidth=1)

        ax.set_xticks(range(len(categories)))
        ax.set_xticklabels(categories)
        ax.set_ylabel(chart_data.get("x_label", "Value ($M)"), fontsize=12)
        ax.set_title(title, fontsize=14, fontweight='bold')
        ax.grid(axis='y', alpha=0.3)

        img_bytes = io.BytesIO()
        plt.savefig(img_bytes, format='png', bbox_inches='tight', dpi=150)
        img_bytes.seek(0)
        plt.close()
        return img_bytes

    def _replace_chart_image(self, slide, chart_data: dict):
        """Remove existing chart images/native charts from slide and add a re-rendered one."""
        # Remove pictures (matplotlib charts) AND native chart shapes
        for shape in list(slide.shapes):
            if shape.shape_type in (13, 3):  # 13=picture, 3=native chart
                shape._element.getparent().remove(shape._element)

        # Re-add appropriate chart type
        chart_type = chart_data.get("chart_type", "bar")
        if chart_type == "waterfall":
            img_bytes = self._render_waterfall_chart(chart_data)
            slide.shapes.add_picture(img_bytes, Inches(1.5), Inches(2), width=Inches(7))
        else:
            self._add_native_bar_chart(slide, chart_data, Inches(1.5), Inches(2.0), Inches(7.0), Inches(4.5))

    # ------------------------------------------------------------------
    # Refinement
    # ------------------------------------------------------------------

    async def refine_presentation(
        self,
        topic: str,
        storyline,
        research,
        length: str,
        feedback,
        iteration: int,
    ) -> str:
        """Load previous PPTX, apply per-slide feedback in-place, save new version."""
        from pptx.util import Pt

        if not self._last_pptx_path:
            # Fallback: regenerate from scratch
            return await self.create_presentation(topic, storyline, research, length)

        prs = Presentation(self._last_pptx_path)

        for fb in feedback:
            slide_idx = fb.slide_index
            if slide_idx >= len(prs.slides):
                continue
            slide = prs.slides[slide_idx]

            # Find title and body textboxes by position
            title_shape = None
            body_shapes = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                try:
                    if shape.top < Inches(1.2):
                        if title_shape is None:
                            title_shape = shape
                    else:
                        body_shapes.append(shape)
                except Exception:
                    body_shapes.append(shape)

            # Replace title
            if fb.new_title and title_shape:
                tf = title_shape.text_frame
                tf.clear()
                para = tf.paragraphs[0]
                para.text = fb.new_title
                para.font.size = Pt(24)
                para.font.bold = True
                para.font.color.rgb = self.primary_color

            # Replace body bullets
            if fb.new_bullets and body_shapes:
                body_shape = body_shapes[0]
                tf = body_shape.text_frame
                tf.clear()
                for i, bullet in enumerate(fb.new_bullets):
                    if i == 0:
                        para = tf.paragraphs[0]
                    else:
                        para = tf.add_paragraph()
                    para.text = bullet
                    para.font.size = Pt(12)
                    para.space_after = Pt(8)

            # Replace chart image (skip title/situation/complication slides 0-2)
            if fb.new_chart_data and slide_idx > 2:
                self._replace_chart_image(slide, fb.new_chart_data)

        os.makedirs("./data/presentations", exist_ok=True)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"presentation_{timestamp}_v{iteration}.pptx"
        filepath = f"./data/presentations/{filename}"
        prs.save(filepath)
        self._last_pptx_path = filepath
        return filepath

    # ------------------------------------------------------------------
    # Standalone chart slides (medium / long)
    # ------------------------------------------------------------------

    def _add_bar_chart_slide(self, prs, storyline: Storyline, research: ResearchResults):
        """Add slide with bar chart and KEY INSIGHT sidebar."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        data = (storyline.slide_data or {}).get("bar_chart", {})
        title = data.get("action_title", "Key Market Drivers")
        self._add_slide_title(slide, title)

        chart_data = {
            "categories": data.get("categories", [f"Factor {i+1}" for i in range(5)]),
            "values": data.get("values", [75, 85, 65, 90, 70]),
            "title": title,
            "x_label": data.get("metric", "Impact Score"),
        }
        self._add_native_bar_chart(slide, chart_data, Inches(0.4), Inches(1.15), Inches(8.6), Inches(5.75))
        top_label, top_value, bullets = self._derive_sidebar_content(chart_data, title)
        self._add_insight_sidebar(slide, title, bullets, top_label, top_value)
        self._add_footer(slide)

    def _add_waterfall_slide(self, prs, storyline: Storyline):
        """Add waterfall chart slide with KEY INSIGHT sidebar."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        data = (storyline.slide_data or {}).get("waterfall", {})
        title = data.get("action_title", "Value Creation Waterfall")
        self._add_slide_title(slide, title)

        chart_data = {
            "categories": data.get("categories", ['Starting', 'Revenue\nGrowth', 'Cost\nReduction', 'Efficiency', 'Ending']),
            "values": data.get("values", [100, 25, 15, 10, 150]),
            "title": title,
            "x_label": data.get("metric", "Value ($M)"),
        }
        img_bytes = self._render_waterfall_chart(chart_data)
        slide.shapes.add_picture(img_bytes, Inches(0.4), Inches(1.15), width=Inches(8.6))
        top_label, top_value, bullets = self._derive_sidebar_content(chart_data, title)
        self._add_insight_sidebar(slide, title, bullets, top_label, top_value)
        self._add_footer(slide)

    def _add_pie_chart_slide(self, prs, storyline: Storyline, research: ResearchResults):
        """Add pie chart slide with KEY INSIGHT sidebar."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        data = (storyline.slide_data or {}).get("pie", {})
        title = data.get("action_title", "Market Segmentation")
        self._add_slide_title(slide, title)

        segments = data.get("categories", ['Segment A', 'Segment B', 'Segment C', 'Segment D', 'Segment E'])
        sizes = data.get("values", [30, 25, 20, 15, 10])
        # Guard mismatched lengths
        n = min(len(segments), len(sizes))
        segments = [_coerce_str(s) for s in segments[:n]]
        sizes = [_coerce_float(v) for v in sizes[:n]]

        cd = ChartData()
        cd.categories = segments
        cd.add_series(data.get("metric", "Market Share"), sizes)
        gf = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.4), Inches(1.15), Inches(8.6), Inches(5.75), cd)
        chart = gf.chart
        chart.has_legend = True
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.show_percentage = True
        plot.data_labels.show_category_name = True

        chart_data = {"categories": segments, "values": sizes, "x_label": data.get("metric", "Market Share")}
        top_label, top_value, bullets = self._derive_sidebar_content(chart_data, title)
        self._add_insight_sidebar(slide, title, bullets, top_label, top_value)
        self._add_footer(slide)

    def _add_tornado_chart_slide(self, prs, storyline: Storyline, research: ResearchResults):
        """Add tornado chart slide with KEY INSIGHT sidebar."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        td = (storyline.slide_data or {}).get("tornado", {})
        title = td.get("action_title", "Sensitivity Analysis")
        self._add_slide_title(slide, title)

        factors = td.get("factors", ['Market Size', 'Pricing', 'Cost Structure', 'Growth Rate', 'Competition'])
        upside = [_coerce_float(v) for v in td.get("upside", [30, 20, 15, 25, 10])]
        downside = [_coerce_float(v) for v in td.get("downside", [-25, -15, -20, -10, -18])]
        # Guard mismatched lengths
        n = min(len(factors), len(upside), len(downside))
        factors, upside, downside = [_coerce_str(f) for f in factors[:n]], upside[:n], downside[:n]

        cd = ChartData()
        cd.categories = factors
        cd.add_series('Upside', upside)
        cd.add_series('Downside', downside)  # already negative values
        gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.4), Inches(1.15), Inches(8.6), Inches(5.75), cd)
        chart = gf.chart
        chart.has_legend = True
        series0 = chart.plots[0].series[0]
        series0.format.fill.solid()
        series0.format.fill.fore_color.rgb = RGBColor(0, 51, 153)   # upside: McKinsey blue
        series1 = chart.plots[0].series[1]
        series1.format.fill.solid()
        series1.format.fill.fore_color.rgb = RGBColor(0, 176, 240)  # downside: light blue

        # Derive sidebar from upside values
        chart_data = {"categories": factors, "values": upside, "x_label": "Impact"}
        top_label, top_value, bullets = self._derive_sidebar_content(chart_data, title)
        self._add_insight_sidebar(slide, title, bullets, top_label, top_value)
        self._add_footer(slide)

    def _add_marimekko_chart_slide(self, prs, storyline: Storyline, research: ResearchResults):
        """Add Marimekko (variable-width stacked bar) chart slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Market Structure")

        fig, ax = plt.subplots(figsize=(8, 4))

        segments = [
            ('Enterprise', 40, [0.6, 0.25, 0.15]),
            ('Mid-Market', 30, [0.4, 0.35, 0.25]),
            ('SMB', 20, [0.3, 0.3, 0.4]),
            ('Consumer', 10, [0.5, 0.3, 0.2]),
        ]
        sub_labels = ['Premium', 'Standard', 'Economy']
        sub_colors = ['#003399', '#00b0f0', '#99ccff']

        x_offset = 0
        for seg_name, width, shares in segments:
            bottom = 0
            for j, (share, color) in enumerate(zip(shares, sub_colors)):
                height = share * 100
                ax.bar(x_offset + width / 2, height, width=width, bottom=bottom,
                       color=color, edgecolor='white', linewidth=1,
                       label=sub_labels[j] if x_offset == 0 else None)
                bottom += height
            ax.text(x_offset + width / 2, -5, seg_name, ha='center', fontsize=9)
            x_offset += width

        ax.set_xlim(0, 100)
        ax.set_ylim(-10, 110)
        ax.set_ylabel('Share (%)', fontsize=12)
        ax.set_title('Market Structure Analysis', fontsize=14, fontweight='bold')
        ax.legend(loc='upper right', fontsize=9)
        ax.set_xticks([])

        img_bytes = io.BytesIO()
        plt.savefig(img_bytes, format='png', bbox_inches='tight', dpi=150)
        img_bytes.seek(0)
        plt.close()

        slide.shapes.add_picture(img_bytes, Inches(0.4), Inches(1.15), width=Inches(8.6))
        self._add_insight_sidebar(slide, "Market Structure",
                                   ["Enterprise leads with 40% market size",
                                    "Premium segment dominates Enterprise tier"])
        self._add_footer(slide)

    def _add_bcg_matrix_slide(self, prs, storyline: Storyline, research: ResearchResults):
        """Add BCG Growth-Share Matrix slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Portfolio Analysis — BCG Growth-Share Matrix")

        fig, ax = plt.subplots(figsize=(7, 4.5))

        ax.axhline(y=10, color='#cccccc', linewidth=1.5, linestyle='--')
        ax.axvline(x=1.0, color='#cccccc', linewidth=1.5, linestyle='--')
        ax.set_facecolor('#f9f9f9')

        quadrant_colors = ['#e8f4fc', '#e8f4fc', '#fef9e7', '#fef9e7']
        for (x0, x1, y0, y1), c in zip(
            [(1.0, 4.0, 10, 30), (0.0, 1.0, 10, 30),
             (1.0, 4.0, 0, 10), (0.0, 1.0, 0, 10)],
            quadrant_colors
        ):
            ax.fill_between([x0, x1], y0, y1, color=c, alpha=0.5)

        ax.text(2.5, 27, 'STARS', ha='center', fontsize=10, fontweight='bold', color='#003399')
        ax.text(0.5, 27, 'QUESTION\nMARKS', ha='center', fontsize=9, color='#666666')
        ax.text(2.5, 2,  'CASH COWS', ha='center', fontsize=10, fontweight='bold', color='#003399')
        ax.text(0.5, 2,  'DOGS', ha='center', fontsize=9, color='#888888')

        units = [
            ('BU-A', 2.4, 22, 1200, '#003399'),
            ('BU-B', 0.6, 18, 500,  '#00b0f0'),
            ('BU-C', 2.8, 5,  2500, '#003399'),
            ('BU-D', 0.4, 3,  400,  '#bbbbbb'),
            ('BU-E', 1.5, 14, 800,  '#0066cc'),
        ]
        for name, x, y, size, color in units:
            ax.scatter(x, y, s=size / 5, color=color, alpha=0.85, edgecolors='white', linewidth=1.5, zorder=5)
            ax.annotate(name, (x, y), textcoords='offset points', xytext=(5, 5),
                        fontsize=9, fontweight='bold', color='#111111')

        ax.set_xlim(0, 4)
        ax.set_ylim(0, 30)
        ax.set_xlabel('Relative Market Share  →  High', fontsize=10)
        ax.set_ylabel('Market Growth Rate (%)  →  High', fontsize=10)
        ax.invert_xaxis()
        ax.grid(False)

        img_bytes = io.BytesIO()
        plt.savefig(img_bytes, format='png', bbox_inches='tight', dpi=150)
        img_bytes.seek(0)
        plt.close()

        slide.shapes.add_picture(img_bytes, Inches(0.4), Inches(1.15), width=Inches(8.6))
        self._add_insight_sidebar(slide, "BCG Portfolio",
                                   ["BU-A and BU-E are Stars — invest for growth",
                                    "BU-C Cash Cow funds portfolio expansion"])
        self._add_footer(slide)

    def _add_priority_matrix_slide(self, prs, storyline: Storyline):
        """Add 2×2 Impact vs. Effort priority matrix slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Prioritization — Impact vs. Effort Matrix")

        fig, ax = plt.subplots(figsize=(7, 4.5))

        fills = [
            (0, 5, 5, 10, '#e3f2fd', 'Quick Wins\n(Do First)'),
            (5, 10, 5, 10, '#fff8e1', 'Strategic\nProjects'),
            (0, 5, 0, 5,  '#f3e5f5', 'Fill-ins\n(Low priority)'),
            (5, 10, 0, 5,  '#fce4ec', 'Avoid /\nDelegate'),
        ]
        for x0, x1, y0, y1, color, label in fills:
            ax.fill_between([x0, x1], y0, y1, color=color, alpha=0.8)
            cx, cy = (x0 + x1) / 2, (y0 + y1) / 2
            ax.text(cx, cy + 1.5, label, ha='center', va='center', fontsize=9,
                    fontweight='bold', color='#444444',
                    bbox=dict(boxstyle='round,pad=0.2', fc='white', alpha=0.6, ec='none'))

        ax.axhline(y=5, color='#888888', linewidth=1.5, linestyle='--')
        ax.axvline(x=5, color='#888888', linewidth=1.5, linestyle='--')

        initiatives = [
            ('Digital Platform', 3, 8.5, 180),
            ('Cost Automation', 2, 7, 120),
            ('M&A Integration', 8, 8, 200),
            ('Analytics BI', 7, 6, 140),
            ('Process Redesign', 3, 3, 90),
            ('Vendor Review', 2, 2, 60),
            ('CRM Upgrade', 8, 2, 80),
        ]
        colors_map = {
            'Digital Platform': '#003399', 'Cost Automation': '#003399',
            'M&A Integration': '#0066cc', 'Analytics BI': '#0066cc',
            'Process Redesign': '#9c27b0', 'Vendor Review': '#9c27b0',
            'CRM Upgrade': '#e53935',
        }
        for name, effort, impact, size in initiatives:
            c = colors_map.get(name, '#003399')
            ax.scatter(effort, impact, s=size, color=c, alpha=0.85,
                       edgecolors='white', linewidth=1.5, zorder=5)
            ax.annotate(name, (effort, impact), textcoords='offset points',
                        xytext=(6, 4), fontsize=8, color='#111111')

        ax.set_xlim(0, 10)
        ax.set_ylim(0, 10)
        ax.set_xlabel('Effort / Complexity  →', fontsize=10)
        ax.set_ylabel('Business Impact  →', fontsize=10)
        ax.set_xticks([])
        ax.set_yticks([])
        ax.text(0.2, 9.5, 'Low Effort', fontsize=8, color='#555555')
        ax.text(8.5, 9.5, 'High Effort', fontsize=8, color='#555555')
        ax.text(0.2, 0.3, 'Low Impact', fontsize=8, color='#555555')
        ax.text(0.2, 9.0, 'High Impact', fontsize=8, color='#555555')

        img_bytes = io.BytesIO()
        plt.savefig(img_bytes, format='png', bbox_inches='tight', dpi=150)
        img_bytes.seek(0)
        plt.close()

        slide.shapes.add_picture(img_bytes, Inches(0.4), Inches(1.15), width=Inches(8.6))
        self._add_insight_sidebar(slide, "Prioritization",
                                   ["Digital Platform and Cost Automation are Quick Wins",
                                    "Focus resources on low-effort, high-impact initiatives first"])
        self._add_footer(slide)

    def _add_value_chain_slide(self, prs, storyline: Storyline):
        """Add Porter Value Chain slide using native PPTX shapes (no image embed)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Value Chain Analysis")

        # ── Primary activities ───────────────────────────────────────────────
        primary = [
            ('Inbound\nLogistics', '#003399'),
            ('Operations', '#0044cc'),
            ('Outbound\nLogistics', '#0055dd'),
            ('Marketing\n& Sales', '#0066ee'),
            ('Service', '#0077ff'),
        ]
        box_w, box_h = Inches(1.5), Inches(1.0)
        top_y = Inches(2.0)
        start_x = Inches(0.5)
        gap = Inches(0.12)

        for i, (label, color) in enumerate(primary):
            x = start_x + i * (box_w + gap)
            shape = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                x, top_y, box_w, box_h
            )
            shape.fill.solid()
            r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
            shape.line.color.rgb = RGBColor(255, 255, 255)
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = label
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

        # Margin box on right
        margin_x = start_x + 5 * (box_w + gap)
        margin_shape = slide.shapes.add_shape(
            1, margin_x, top_y, Inches(1.1), box_h
        )
        margin_shape.fill.solid()
        margin_shape.fill.fore_color.rgb = RGBColor(0, 51, 153)
        margin_shape.line.color.rgb = RGBColor(255, 255, 255)
        tf = margin_shape.text_frame
        p = tf.paragraphs[0]
        p.text = 'Margin'
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # ── Support activities ───────────────────────────────────────────────
        support = [
            ('Firm Infrastructure', '#4a6fa5'),
            ('Human Resource Management', '#5d7db5'),
            ('Technology Development', '#7090c4'),
            ('Procurement', '#839fd3'),
        ]
        sup_h = Inches(0.55)
        sup_w = Inches(8.1)  # full width
        sup_x = start_x
        for i, (label, color) in enumerate(support):
            y = top_y + box_h + Inches(0.12) + i * (sup_h + Inches(0.06))
            s = slide.shapes.add_shape(1, sup_x, y, sup_w, sup_h)
            s.fill.solid()
            r, g, b = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
            s.fill.fore_color.rgb = RGBColor(r, g, b)
            s.line.color.rgb = RGBColor(255, 255, 255)
            tf = s.text_frame
            p = tf.paragraphs[0]
            p.text = label
            p.alignment = PP_ALIGN.LEFT
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(255, 255, 255)

        # Label
        lbl = slide.shapes.add_textbox(start_x, Inches(1.72), Inches(4), Inches(0.25))
        lp = lbl.text_frame.paragraphs[0]
        lp.text = 'Primary Activities'
        lp.font.size = Pt(9)
        lp.font.bold = True
        lp.font.color.rgb = RGBColor(0, 51, 153)

        lbl2 = slide.shapes.add_textbox(start_x, top_y + box_h + Inches(0.05), Inches(2), Inches(0.22))
        lp2 = lbl2.text_frame.paragraphs[0]
        lp2.text = 'Support Activities'
        lp2.font.size = Pt(9)
        lp2.font.bold = True
        lp2.font.color.rgb = RGBColor(0, 51, 153)

    def _add_heatmap_slide(self, prs, storyline: Storyline, research: ResearchResults):
        """Add competitive landscape heatmap slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Competitive Landscape — Capability Heatmap")

        fig, ax = plt.subplots(figsize=(8, 4))

        capabilities = ['Digital', 'Operations', 'Talent', 'Innovation', 'Customer\nExp.', 'Cost\nEfficiency']
        competitors = ['Our Co.', 'Competitor A', 'Competitor B', 'Competitor C', 'Competitor D']

        np.random.seed(42)
        data = np.array([
            [9, 7, 8, 8, 9, 7],
            [6, 8, 5, 6, 7, 9],
            [7, 6, 7, 5, 6, 8],
            [5, 7, 6, 7, 5, 6],
            [4, 5, 4, 6, 4, 7],
        ], dtype=float)

        cmap = plt.cm.RdYlGn
        im = ax.imshow(data, cmap=cmap, vmin=1, vmax=10, aspect='auto')

        ax.set_xticks(range(len(capabilities)))
        ax.set_xticklabels(capabilities, fontsize=10)
        ax.set_yticks(range(len(competitors)))
        ax.set_yticklabels(competitors, fontsize=10)

        for i in range(len(competitors)):
            for j in range(len(capabilities)):
                val = data[i, j]
                text_color = 'white' if val < 4 or val > 8 else 'black'
                ax.text(j, i, f'{val:.0f}', ha='center', va='center',
                        fontsize=11, fontweight='bold', color=text_color)

        for j in range(len(capabilities)):
            ax.add_patch(plt.Rectangle((j - 0.5, -0.5), 1, 1,
                                        fill=False, edgecolor='#003399', linewidth=2.5))

        plt.colorbar(im, ax=ax, label='Score (1–10)', shrink=0.8)
        ax.set_title('Competitive Capability Assessment', fontsize=13, fontweight='bold', pad=10)

        img_bytes = io.BytesIO()
        plt.savefig(img_bytes, format='png', bbox_inches='tight', dpi=150)
        img_bytes.seek(0)
        plt.close()

        slide.shapes.add_picture(img_bytes, Inches(0.4), Inches(1.15), width=Inches(8.6))
        self._add_insight_sidebar(slide, "Competitive Landscape",
                                   ["Our Co. leads in Digital and Customer Experience",
                                    "Competitor A strongest in Operations and Cost Efficiency"])
        self._add_footer(slide)

    # ------------------------------------------------------------------
    # Recommendations & Sources
    # ------------------------------------------------------------------

    def _add_recommendations(self, prs, storyline: Storyline):
        """Add recommendations slide with numbered oval badges."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        answer = storyline.scqa.answer
        title_text = (answer[:70] + "...") if len(answer) > 70 else answer
        self._add_slide_title(slide, title_text)

        items = (
            storyline.recommendation_items
            if storyline.recommendation_items
            else [s.strip() for s in storyline.scqa.answer.split(". ") if s.strip()][:5]
        )

        y = Inches(1.15)
        item_h = Inches(0.85)
        for i, item in enumerate(items):
            # Blue oval badge
            oval = slide.shapes.add_shape(9, Inches(0.4), y + Inches(0.1),
                                           Inches(0.55), Inches(0.55))
            oval.fill.solid()
            oval.fill.fore_color.rgb = RGBColor(0, 51, 153)
            oval.line.fill.background()

            num_box = slide.shapes.add_textbox(Inches(0.4), y + Inches(0.08),
                                                Inches(0.55), Inches(0.55))
            np_ = num_box.text_frame.paragraphs[0]
            np_.text = str(i + 1)
            np_.font.size = Pt(14)
            np_.font.bold = True
            np_.font.color.rgb = RGBColor(255, 255, 255)
            np_.alignment = PP_ALIGN.CENTER

            # Item text box — 12.0" wide on widescreen
            item_box = slide.shapes.add_textbox(Inches(1.1), y, Inches(12.0), item_h)
            tf = item_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = self._strip_markdown(item)
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 153)

            y += item_h + Inches(0.1)

        self._add_footer(slide)

    def _add_sources(self, prs, research: ResearchResults):
        """Add sources/references slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_title(slide, "Sources")

        content_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.0), Inches(5.0))
        tf = content_box.text_frame
        tf.word_wrap = True

        # Collect unique sources
        sources_seen = set()
        source_num = 1

        for hyp_evidence in research.hypotheses_evidence:
            for evidence in hyp_evidence.evidence[:3]:  # Top 3 per hypothesis
                if evidence.url not in sources_seen:
                    sources_seen.add(evidence.url)
                    p = tf.add_paragraph() if source_num > 1 else tf.paragraphs[0]
                    p.text = f"[{source_num}] {evidence.source} - {evidence.url}"
                    p.font.size = Pt(8)
                    p.space_after = Pt(6)
                    source_num += 1

                    if source_num > 15:  # Limit sources
                        break
