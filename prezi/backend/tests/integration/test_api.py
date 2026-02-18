"""Integration tests for FastAPI endpoints via httpx."""

import pytest
from unittest.mock import patch
from app.database import Job


class TestRootEndpoints:
    async def test_root(self, test_client):
        """GET / → service info."""
        resp = await test_client.get("/")
        assert resp.status_code == 200
        data = resp.json()
        assert data["service"] == "Prezi AI"
        assert data["version"] == "1.0.0"
        assert data["status"] == "running"

    async def test_health(self, test_client):
        """GET /health → healthy."""
        resp = await test_client.get("/health")
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "healthy"


class TestProvidersEndpoint:
    async def test_providers(self, test_client):
        """GET /api/providers → mock research always present."""
        resp = await test_client.get("/api/providers")
        assert resp.status_code == 200
        data = resp.json()

        research_ids = [p["id"] for p in data["research_providers"]]
        assert "mock" in research_ids

        mock = next(p for p in data["research_providers"] if p["id"] == "mock")
        assert mock["available"] is True


class TestGenerateEndpoint:
    async def test_generate_invalid_provider(self, test_client):
        """POST /api/generate with unavailable LLM → 400."""
        resp = await test_client.post(
            "/api/generate",
            json={
                "topic": "Cloud computing strategy for enterprise clients",
                "length": "short",
                "llm_provider": "nonexistent_provider",
                "research_provider": "mock",
            },
        )
        assert resp.status_code == 400

    async def test_generate_topic_too_short(self, test_client):
        """POST with short topic → 422."""
        resp = await test_client.post(
            "/api/generate",
            json={
                "topic": "Short",
                "length": "short",
                "llm_provider": "claude",
                "research_provider": "mock",
            },
        )
        assert resp.status_code == 422


class TestStatusEndpoint:
    async def test_status_not_found(self, test_client):
        """GET /api/status/fake → 404."""
        resp = await test_client.get("/api/status/nonexistent-job-id")
        assert resp.status_code == 404


class TestDownloadEndpoint:
    async def test_download_not_found(self, test_client):
        """GET /api/download/fake → 404."""
        resp = await test_client.get("/api/download/nonexistent-job-id")
        assert resp.status_code == 404

    async def test_download_incomplete(self, test_client, db_engine):
        """Incomplete job → 400."""
        from sqlalchemy.orm import sessionmaker

        Session = sessionmaker(bind=db_engine)
        session = Session()
        job = Job(
            id="incomplete-dl",
            topic="Testing incomplete download scenario here",
            length="short",
            llm_provider="claude",
            research_provider="mock",
            status="storyline",
            progress=20,
            message="Generating...",
        )
        session.add(job)
        session.commit()
        session.close()

        resp = await test_client.get("/api/download/incomplete-dl")
        assert resp.status_code == 400


class TestRetryEndpoint:
    async def test_retry_not_found(self, test_client):
        """POST /api/retry/fake → 404."""
        resp = await test_client.post("/api/retry/nonexistent-job-id")
        assert resp.status_code == 404

    async def test_retry_non_failed_job(self, test_client, db_engine):
        """Retry a non-failed job → 400."""
        from sqlalchemy.orm import sessionmaker

        Session = sessionmaker(bind=db_engine)
        session = Session()
        job = Job(
            id="not-failed",
            topic="Testing retry on a non-failed job here",
            length="short",
            llm_provider="claude",
            research_provider="mock",
            status="completed",
            progress=100,
            message="Done",
        )
        session.add(job)
        session.commit()
        session.close()

        resp = await test_client.post("/api/retry/not-failed")
        assert resp.status_code == 400
        assert "Can only retry failed jobs" in resp.json()["detail"]

    @patch("app.api.endpoints.generate_presentation_background")
    async def test_retry_failed_job(self, mock_bg, test_client, db_engine):
        """Retry a failed job → resets to queued, returns 200."""
        from sqlalchemy.orm import sessionmaker

        Session = sessionmaker(bind=db_engine)
        session = Session()
        job = Job(
            id="failed-job",
            topic="Testing retry on a failed generation job",
            length="medium",
            llm_provider="claude",
            research_provider="mock",
            status="failed",
            progress=30,
            message="Failed: some error",
            error="some error",
        )
        session.add(job)
        session.commit()
        session.close()

        resp = await test_client.post("/api/retry/failed-job")
        assert resp.status_code == 200
        assert resp.json()["job_id"] == "failed-job"

        # Verify job was reset
        session = Session()
        updated = session.query(Job).filter(Job.id == "failed-job").first()
        assert updated.status == "queued"
        assert updated.progress == 0
        assert updated.error is None
        session.close()

        mock_bg.assert_called_once()


class TestJobsEndpoint:
    async def test_jobs_empty(self, test_client):
        """GET /api/jobs with empty DB → empty list."""
        resp = await test_client.get("/api/jobs")
        assert resp.status_code == 200
        data = resp.json()
        assert data["jobs"] == []
        assert data["total"] == 0
        assert data["page"] == 1

    async def test_jobs_pagination(self, test_client, db_engine):
        """Pagination returns correct page/per_page."""
        from sqlalchemy.orm import sessionmaker

        Session = sessionmaker(bind=db_engine)
        session = Session()
        for i in range(5):
            session.add(
                Job(
                    id=f"page-job-{i}",
                    topic=f"Testing pagination with job number {i}",
                    length="short",
                    llm_provider="claude",
                    research_provider="mock",
                    status="completed",
                    progress=100,
                    message="Done",
                )
            )
        session.commit()
        session.close()

        resp = await test_client.get("/api/jobs?page=1&per_page=2")
        assert resp.status_code == 200
        data = resp.json()
        assert len(data["jobs"]) == 2
        assert data["total"] == 5
        assert data["page"] == 1
        assert data["per_page"] == 2

    async def test_jobs_with_quality_score(self, test_client, db_engine):
        """Completed job includes overall quality score."""
        from sqlalchemy.orm import sessionmaker

        Session = sessionmaker(bind=db_engine)
        session = Session()
        session.add(
            Job(
                id="quality-job",
                topic="Testing quality score extraction from job",
                length="medium",
                llm_provider="claude",
                research_provider="mock",
                status="completed",
                progress=100,
                message="Done",
                quality_score={"overall_score": 85, "slide_logic": 90},
            )
        )
        session.commit()
        session.close()

        resp = await test_client.get("/api/jobs")
        assert resp.status_code == 200
        data = resp.json()
        job = data["jobs"][0]
        assert job["quality_score_overall"] == 85


class TestPdfEndpoint:
    async def test_pdf_not_found(self, test_client):
        """GET /api/download/fake/pdf → 404."""
        resp = await test_client.get("/api/download/nonexistent-job-id/pdf")
        assert resp.status_code == 404

    async def test_pdf_incomplete_job(self, test_client, db_engine):
        """Incomplete job → 400."""
        from sqlalchemy.orm import sessionmaker

        Session = sessionmaker(bind=db_engine)
        session = Session()
        session.add(
            Job(
                id="incomplete-pdf",
                topic="Testing PDF download on incomplete job",
                length="short",
                llm_provider="claude",
                research_provider="mock",
                status="slides",
                progress=65,
                message="Creating slides...",
            )
        )
        session.commit()
        session.close()

        resp = await test_client.get("/api/download/incomplete-pdf/pdf")
        assert resp.status_code == 400

    @patch("app.api.endpoints.shutil.which", return_value=None)
    async def test_pdf_no_libreoffice(self, mock_which, test_client, db_engine):
        """No LibreOffice installed → 503."""
        from sqlalchemy.orm import sessionmaker

        Session = sessionmaker(bind=db_engine)
        session = Session()
        session.add(
            Job(
                id="no-libre",
                topic="Testing PDF with no LibreOffice installed",
                length="short",
                llm_provider="claude",
                research_provider="mock",
                status="completed",
                progress=100,
                message="Done",
                pptx_path="/tmp/fake.pptx",
            )
        )
        session.commit()
        session.close()

        resp = await test_client.get("/api/download/no-libre/pdf")
        assert resp.status_code == 503
        assert "LibreOffice" in resp.json()["detail"]


class TestTemplateEndpoints:
    async def test_template_list_has_default(self, test_client):
        """GET /api/templates always includes 'default'."""
        resp = await test_client.get("/api/templates")
        assert resp.status_code == 200
        data = resp.json()
        ids = [t["id"] for t in data["templates"]]
        assert "default" in ids
        default = next(t for t in data["templates"] if t["id"] == "default")
        assert default["name"] == "McKinsey Classic"

    async def test_template_upload_valid(self, test_client):
        """Upload a valid .pptx file → 200."""
        from pptx import Presentation as PptxPresentation
        import io

        # Create a minimal valid PPTX in memory
        prs = PptxPresentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        resp = await test_client.post(
            "/api/templates/upload",
            files={"file": ("test_template.pptx", buf, "application/octet-stream")},
            data={"name": "Test Template"},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["name"] == "Test Template"
        assert data["filename"] == "test_template.pptx"
        assert data["id"]

    async def test_template_upload_invalid_extension(self, test_client):
        """Upload a non-.pptx file → 400."""
        import io

        resp = await test_client.post(
            "/api/templates/upload",
            files={"file": ("bad.txt", io.BytesIO(b"not a pptx"), "text/plain")},
            data={"name": "Bad File"},
        )
        assert resp.status_code == 400
        assert "pptx" in resp.json()["detail"].lower()

    async def test_template_delete(self, test_client):
        """Upload then delete a template."""
        from pptx import Presentation as PptxPresentation
        import io

        prs = PptxPresentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        resp = await test_client.post(
            "/api/templates/upload",
            files={"file": ("del_template.pptx", buf, "application/octet-stream")},
            data={"name": "To Delete"},
        )
        template_id = resp.json()["id"]

        # Delete it
        resp = await test_client.delete(f"/api/templates/{template_id}")
        assert resp.status_code == 200

        # Verify it's gone from list
        resp = await test_client.get("/api/templates")
        ids = [t["id"] for t in resp.json()["templates"]]
        assert template_id not in ids

    async def test_template_delete_default_blocked(self, test_client):
        """Cannot delete the default template → 400."""
        resp = await test_client.delete("/api/templates/default")
        assert resp.status_code == 400

    async def test_generate_accepts_template_id(self, test_client):
        """Generate request with template_id field is accepted (validation only)."""
        resp = await test_client.post(
            "/api/generate",
            json={
                "topic": "Cloud computing strategy for enterprise clients",
                "length": "short",
                "llm_provider": "nonexistent_provider",
                "research_provider": "mock",
                "template_id": "default",
            },
        )
        # 400 because LLM provider doesn't exist, but template_id is accepted
        assert resp.status_code == 400
        assert "LLM provider" in resp.json()["detail"]


class TestResultEndpoint:
    async def test_result_not_found(self, test_client):
        """GET /api/result/fake → 404."""
        resp = await test_client.get("/api/result/nonexistent-job-id")
        assert resp.status_code == 404

    async def test_result_incomplete(self, test_client, db_engine):
        """Incomplete job → 400."""
        from sqlalchemy.orm import sessionmaker

        Session = sessionmaker(bind=db_engine)
        session = Session()
        job = Job(
            id="incomplete-res",
            topic="Testing incomplete result scenario here",
            length="medium",
            llm_provider="claude",
            research_provider="mock",
            status="researching",
            progress=40,
            message="Researching...",
        )
        session.add(job)
        session.commit()
        session.close()

        resp = await test_client.get("/api/result/incomplete-res")
        assert resp.status_code == 400
