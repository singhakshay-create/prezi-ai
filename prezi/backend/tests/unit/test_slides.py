"""Tests for SlideGenerator agent."""

import os
import pytest
from pptx import Presentation
from pptx.util import Inches

from app.agents.slides import SlideGenerator, LAYOUT_REGISTRY
from app.models import Storyline, ResearchResults, Hypothesis, SCQAFramework, HypothesisEvidence, SearchResult


class TestSlideGenerator:
    """Tests for SlideGenerator.create_presentation()."""

    async def test_short_slide_count(self, sample_storyline, sample_research_results):
        """Short deck: title + situation + complication + 3 hypothesis slides + recommendations + sources = 8."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "short"
        )
        try:
            prs = Presentation(path)
            assert len(prs.slides) == 8
        finally:
            os.remove(path)

    async def test_medium_slide_count(self, sample_storyline, sample_research_results):
        """Medium deck: short (8) + bar + waterfall + pie + tornado = 12."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "medium"
        )
        try:
            prs = Presentation(path)
            assert len(prs.slides) == 12
        finally:
            os.remove(path)

    async def test_long_slide_count(self, sample_storyline, sample_research_results):
        """Long deck: medium (12) + marimekko + BCG + priority + value chain + heatmap = 17."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "long"
        )
        try:
            prs = Presentation(path)
            assert len(prs.slides) == 17
        finally:
            os.remove(path)

    async def test_pptx_valid(self, sample_storyline, sample_research_results):
        """Output file opens with Presentation() without error."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "short"
        )
        try:
            prs = Presentation(path)
            assert prs is not None
        finally:
            os.remove(path)

    async def test_title_contains_topic(self, sample_storyline, sample_research_results):
        """First slide text includes the topic string."""
        gen = SlideGenerator()
        topic = "Cloud Strategy for Fortune 500"
        path = await gen.create_presentation(
            topic, sample_storyline, sample_research_results, "short"
        )
        try:
            prs = Presentation(path)
            first_slide = prs.slides[0]
            texts = [shape.text_frame.text for shape in first_slide.shapes if shape.has_text_frame]
            combined = " ".join(texts)
            assert topic in combined
        finally:
            os.remove(path)

    async def test_slide_dimensions(self, sample_storyline, sample_research_results):
        """Slide dimensions are 13.333" x 7.5" (widescreen 16:9)."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "short"
        )
        try:
            prs = Presentation(path)
            assert prs.slide_width == Inches(13.333)
            assert prs.slide_height == Inches(7.5)
        finally:
            os.remove(path)

    async def test_output_path_pattern(self, sample_storyline, sample_research_results):
        """File saved to ./data/presentations/presentation_*.pptx."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "short"
        )
        try:
            assert path.startswith("./data/presentations/presentation_")
            assert path.endswith(".pptx")
            assert os.path.isfile(path)
        finally:
            os.remove(path)

    async def test_medium_chart_counts(self, sample_storyline, sample_research_results):
        """Medium deck: 1 matplotlib PNG (waterfall) + 6 native charts (3 hyp + bar + pie + tornado)."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "medium"
        )
        try:
            prs = Presentation(path)
            picture_count = sum(
                1 for slide in prs.slides for shape in slide.shapes
                if shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE (matplotlib PNGs)
            )
            chart_count = sum(
                1 for slide in prs.slides for shape in slide.shapes
                if shape.shape_type == 3  # native pptx chart objects
            )
            assert picture_count == 1, f"Expected 1 waterfall PNG, got {picture_count}"
            assert chart_count == 6, f"Expected 6 native charts, got {chart_count}"
        finally:
            os.remove(path)

    async def test_long_chart_counts(self, sample_storyline, sample_research_results):
        """Long deck: 5 matplotlib PNGs (waterfall + marimekko + BCG + priority + heatmap) + 6 native charts.
        Value chain uses native PPTX rectangles (not pictures or chart objects)."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "long"
        )
        try:
            prs = Presentation(path)
            picture_count = sum(
                1 for slide in prs.slides for shape in slide.shapes
                if shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE (matplotlib PNGs)
            )
            chart_count = sum(
                1 for slide in prs.slides for shape in slide.shapes
                if shape.shape_type == 3  # native pptx chart objects
            )
            assert picture_count == 5, f"Expected 5 matplotlib PNGs, got {picture_count}"
            assert chart_count == 6, f"Expected 6 native charts, got {chart_count}"
        finally:
            os.remove(path)

    async def test_value_chain_has_native_shapes(self, sample_storyline, sample_research_results):
        """Value chain slide uses native PPTX rectangles, no embedded images.
        Slide order (long): 0=title 1=situation 2=complication 3-5=hypotheses 6=bar 7=waterfall
        8=pie 9=tornado 10=marimekko 11=BCG 12=priority 13=value-chain 14=heatmap 15=recs 16=sources"""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "long"
        )
        try:
            prs = Presentation(path)
            assert len(prs.slides) == 17
            value_chain_slide = prs.slides[13]
            pictures = [s for s in value_chain_slide.shapes if s.shape_type == 13]
            rectangles = [s for s in value_chain_slide.shapes if s.shape_type == 1]
            assert len(pictures) == 0
            assert len(rectangles) >= 9  # 5 primary + 4 support activity boxes
        finally:
            os.remove(path)

    async def test_slides_with_template_path(self, sample_storyline, sample_research_results):
        """SlideGenerator with a template_path produces a valid PPTX."""
        # Create a minimal template
        template_prs = Presentation()
        template_prs.slide_width = Inches(10)
        template_prs.slide_height = Inches(7.5)
        os.makedirs("./data/templates", exist_ok=True)
        template_path = "./data/templates/test_template.pptx"
        template_prs.save(template_path)

        try:
            gen = SlideGenerator(template_path=template_path)
            path = await gen.create_presentation(
                "Cloud Strategy", sample_storyline, sample_research_results, "short"
            )
            try:
                prs = Presentation(path)
                assert prs is not None
                assert len(prs.slides) == 8
            finally:
                os.remove(path)
        finally:
            os.remove(template_path)

    async def test_refine_presentation_replaces_title(
        self, sample_storyline, sample_research_results
    ):
        """refine_presentation with new_title changes the title on the target slide."""
        from app.models import SlideFeedback

        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "short"
        )

        feedback = [
            SlideFeedback(
                slide_index=1,
                new_title="Hybrid Cloud Adoption Grows 2× Faster",
                new_bullets=None,
                new_chart_data=None,
                issues_addressed=["weak_title"],
            )
        ]
        try:
            refined_path = await gen.refine_presentation(
                "Cloud Strategy", sample_storyline, sample_research_results, "short",
                feedback, 1
            )
            try:
                prs = Presentation(refined_path)
                slide = prs.slides[1]
                texts = [
                    shape.text_frame.text
                    for shape in slide.shapes
                    if shape.has_text_frame
                ]
                combined = " ".join(texts)
                assert "Hybrid Cloud Adoption Grows 2× Faster" in combined
            finally:
                if os.path.isfile(refined_path):
                    os.remove(refined_path)
        finally:
            if os.path.isfile(path):
                os.remove(path)

    async def test_refine_presentation_replaces_chart(
        self, sample_storyline, sample_research_results
    ):
        """refine_presentation with new_chart_data keeps a picture shape on the slide."""
        from app.models import SlideFeedback

        gen = SlideGenerator()
        # Use medium deck so slide index 3 has a hypothesis bar chart
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "medium"
        )

        feedback = [
            SlideFeedback(
                slide_index=3,
                new_title=None,
                new_bullets=None,
                new_chart_data={
                    "chart_type": "bar",
                    "categories": ["Hybrid Cloud", "Public Cloud", "On-Premises"],
                    "values": [85, 75, 45],
                    "title": "Adoption by Model",
                    "x_label": "Score",
                },
                issues_addressed=["placeholder_data"],
            )
        ]
        try:
            refined_path = await gen.refine_presentation(
                "Cloud Strategy", sample_storyline, sample_research_results, "medium",
                feedback, 1
            )
            try:
                prs = Presentation(refined_path)
                slide = prs.slides[3]
                chart_or_pic = [s for s in slide.shapes if s.shape_type in (3, 13)]
                assert len(chart_or_pic) >= 1
            finally:
                if os.path.isfile(refined_path):
                    os.remove(refined_path)
        finally:
            if os.path.isfile(path):
                os.remove(path)

    def test_render_bar_chart_with_real_data(self):
        """_render_bar_chart with explicit data returns non-empty BytesIO."""
        import io as _io

        gen = SlideGenerator()
        chart_data = {
            "categories": ["Hybrid Cloud", "Public Cloud", "On-Premises"],
            "values": [85, 75, 45],
            "title": "Adoption by Model",
            "x_label": "Score (%)",
        }
        result = gen._render_bar_chart(chart_data)
        assert isinstance(result, _io.BytesIO)
        result.seek(0)
        assert len(result.read()) > 0

    def test_render_bar_chart_more_categories_than_values(self):
        """_render_bar_chart does not crash when categories outnumber values."""
        import io as _io

        gen = SlideGenerator()
        result = gen._render_bar_chart({
            "categories": ["A", "B", "C", "D", "E"],
            "values": [10, 20],
        })
        result.seek(0)
        assert len(result.read()) > 0

    def test_render_bar_chart_more_values_than_categories(self):
        """_render_bar_chart does not crash when values outnumber categories."""
        import io as _io

        gen = SlideGenerator()
        result = gen._render_bar_chart({
            "categories": ["A", "B"],
            "values": [10, 20, 30, 40, 50],
        })
        result.seek(0)
        assert len(result.read()) > 0

    def test_render_waterfall_chart_mismatched_lengths(self):
        """_render_waterfall_chart does not crash when categories and values differ in length."""
        import io as _io

        gen = SlideGenerator()
        result = gen._render_waterfall_chart({
            "categories": ["Start", "Growth", "Cost", "Efficiency", "End"],
            "values": [100, 50],
            "title": "Waterfall",
        })
        result.seek(0)
        assert len(result.read()) > 0

    def test_render_bar_chart_values_as_dicts(self):
        """_render_bar_chart does not crash when LLM returns values as dicts instead of numbers."""
        import io as _io

        gen = SlideGenerator()
        result = gen._render_bar_chart({
            "categories": ["Hybrid Cloud", "Public Cloud", "On-Premises"],
            "values": [{"value": 85}, {"value": 75}, {"value": 45}],
            "title": "Adoption by Model",
            "x_label": "Score",
        })
        result.seek(0)
        assert len(result.read()) > 0

    def test_render_bar_chart_categories_as_dicts(self):
        """_render_bar_chart does not crash when LLM returns categories as dicts instead of strings."""
        import io as _io

        gen = SlideGenerator()
        result = gen._render_bar_chart({
            "categories": [{"label": "Hybrid Cloud"}, {"label": "Public Cloud"}],
            "values": [85, 75],
            "title": "Adoption by Model",
            "x_label": "Score",
        })
        result.seek(0)
        assert len(result.read()) > 0

    def test_render_waterfall_chart_values_as_dicts(self):
        """_render_waterfall_chart does not crash when LLM returns values as dicts."""
        import io as _io

        gen = SlideGenerator()
        result = gen._render_waterfall_chart({
            "categories": ["Start", "Growth", "End"],
            "values": [{"v": 100}, {"v": 30}, {"v": 130}],
            "title": "Waterfall",
        })
        result.seek(0)
        assert len(result.read()) > 0

    # ------------------------------------------------------------------
    # New tests for consulting-quality redesign
    # ------------------------------------------------------------------

    async def test_short_hypothesis_slides_have_charts(
        self, sample_storyline, sample_research_results
    ):
        """Slides 3, 4, 5 (one per hypothesis) each have at least 1 native chart shape."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "short"
        )
        try:
            prs = Presentation(path)
            for idx in [3, 4, 5]:
                slide = prs.slides[idx]
                chart_shapes = [s for s in slide.shapes if s.shape_type == 3]
                assert len(chart_shapes) >= 1, f"Slide {idx} missing native chart"
        finally:
            os.remove(path)

    async def test_situation_slide_uses_action_title(
        self, sample_storyline, sample_research_results
    ):
        """Slide 1 (situation) contains the situation_title text."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "short"
        )
        try:
            prs = Presentation(path)
            slide = prs.slides[1]
            texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
            combined = " ".join(texts)
            assert sample_storyline.scqa.situation_title in combined
        finally:
            os.remove(path)

    async def test_complication_slide_uses_action_title(
        self, sample_storyline, sample_research_results
    ):
        """Slide 2 (complication) contains the complication_title text."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "short"
        )
        try:
            prs = Presentation(path)
            slide = prs.slides[2]
            texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
            combined = " ".join(texts)
            assert sample_storyline.scqa.complication_title in combined
        finally:
            os.remove(path)

    async def test_hypothesis_slide_title_is_action_title(
        self, sample_storyline, sample_research_results
    ):
        """Slide 3 (first hypothesis) contains the action_title text."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "short"
        )
        try:
            prs = Presentation(path)
            slide = prs.slides[3]
            texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
            combined = " ".join(texts)
            assert sample_storyline.hypotheses[0].action_title in combined
        finally:
            os.remove(path)

    async def test_bar_chart_slide_uses_slide_data_title(
        self, sample_storyline, sample_research_results
    ):
        """Medium deck: slide 6 (bar chart) title matches slide_data['bar_chart']['action_title']."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "medium"
        )
        try:
            prs = Presentation(path)
            # Slide 6: title + situation + complication + 3 hyp = 6 slides before bar chart
            slide = prs.slides[6]
            texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
            combined = " ".join(texts)
            assert sample_storyline.slide_data["bar_chart"]["action_title"] in combined
        finally:
            os.remove(path)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_hypothesis(chart_type: str, hyp_id: int = 1) -> Hypothesis:
    return Hypothesis(
        id=hyp_id,
        text=f"Hypothesis {hyp_id}",
        testable_claim=f"Claim {hyp_id}",
        action_title=f"Action title for hypothesis {hyp_id}",
        chart_hint={
            "type": chart_type,
            "categories": ["Phase 1", "Phase 2", "Phase 3"],
            "values": [10, 25, 40],
            "metric": "Score",
        },
    )


def _empty_research() -> ResearchResults:
    return ResearchResults(hypotheses_evidence=[], total_sources=0)


def _single_hyp_storyline(chart_type: str) -> Storyline:
    return Storyline(
        scqa=SCQAFramework(
            situation="Market growing rapidly.",
            complication="Competitors are catching up.",
            question="How do we maintain our lead?",
            answer="Three strategic actions.",
            situation_title="Market grows 25% CAGR",
            complication_title="Competitive intensity rising",
        ),
        governing_thought="Act decisively in next 12 months.",
        key_line="Three pillars drive sustainable advantage.",
        hypotheses=[_make_hypothesis(chart_type)],
        recommendation_items=["Action 1: do X", "Action 2: do Y"],
        slide_data={},
    )


# ---------------------------------------------------------------------------
# Layout registry
# ---------------------------------------------------------------------------

class TestLayoutRegistry:

    def test_registry_has_expected_types(self):
        for t in ["bar", "waterfall", "pie", "tornado", "timeline", "three_kpi", "two_column", "default"]:
            assert t in LAYOUT_REGISTRY

    def test_bar_maps_to_chart_evidence(self):
        assert LAYOUT_REGISTRY["bar"] == "chart_evidence"

    def test_timeline_maps_to_timeline(self):
        assert LAYOUT_REGISTRY["timeline"] == "timeline"

    def test_three_kpi_maps_to_three_kpi(self):
        assert LAYOUT_REGISTRY["three_kpi"] == "three_kpi"

    def test_two_column_maps_to_two_column(self):
        assert LAYOUT_REGISTRY["two_column"] == "two_column"


# ---------------------------------------------------------------------------
# New layout slides
# ---------------------------------------------------------------------------

class TestNewLayouts:

    async def test_timeline_slide_produces_valid_pptx(self):
        """A hypothesis with chart_hint.type='timeline' generates a valid PPTX."""
        gen = SlideGenerator()
        storyline = _single_hyp_storyline("timeline")
        path = await gen.create_presentation("Cloud Roadmap", storyline, _empty_research(), "short")
        try:
            prs = Presentation(path)
            assert len(prs.slides) >= 1
        finally:
            os.remove(path)

    async def test_three_kpi_slide_produces_valid_pptx(self):
        """A hypothesis with chart_hint.type='three_kpi' generates a valid PPTX."""
        gen = SlideGenerator()
        storyline = _single_hyp_storyline("three_kpi")
        path = await gen.create_presentation("KPI Overview", storyline, _empty_research(), "short")
        try:
            prs = Presentation(path)
            assert len(prs.slides) >= 1
        finally:
            os.remove(path)

    async def test_two_column_slide_produces_valid_pptx(self):
        """A hypothesis with chart_hint.type='two_column' generates a valid PPTX."""
        gen = SlideGenerator()
        storyline = _single_hyp_storyline("two_column")
        path = await gen.create_presentation("Analysis", storyline, _empty_research(), "short")
        try:
            prs = Presentation(path)
            assert len(prs.slides) >= 1
        finally:
            os.remove(path)

    async def test_timeline_slide_uses_action_title(self):
        """Timeline slide title contains the hypothesis action_title."""
        gen = SlideGenerator()
        storyline = _single_hyp_storyline("timeline")
        path = await gen.create_presentation("Cloud Roadmap", storyline, _empty_research(), "short")
        try:
            prs = Presentation(path)
            # Slide 3 is the hypothesis slide (0=title, 1=situation, 2=complication)
            slide = prs.slides[3]
            texts = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
            assert "Action title for hypothesis 1" in texts
        finally:
            os.remove(path)

    async def test_three_kpi_shows_categories(self):
        """Three KPI slide contains category labels from chart_hint."""
        gen = SlideGenerator()
        storyline = _single_hyp_storyline("three_kpi")
        path = await gen.create_presentation("KPIs", storyline, _empty_research(), "short")
        try:
            prs = Presentation(path)
            slide = prs.slides[3]
            texts = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
            assert "Phase 1" in texts
        finally:
            os.remove(path)

    async def test_unknown_chart_type_falls_back_to_bar(self):
        """Unrecognised chart_hint.type falls back to chart_evidence layout (bar chart)."""
        gen = SlideGenerator()
        storyline = _single_hyp_storyline("marimekko_custom")  # not in registry
        path = await gen.create_presentation("Analysis", storyline, _empty_research(), "short")
        try:
            prs = Presentation(path)
            slide = prs.slides[3]
            # Fallback: bar chart present (native chart shape type 3)
            charts = [s for s in slide.shapes if s.shape_type == 3]
            assert len(charts) >= 1
        finally:
            os.remove(path)

    async def test_timeline_with_evidence(self):
        """Timeline slide renders correctly when evidence is provided."""
        gen = SlideGenerator()
        storyline = _single_hyp_storyline("timeline")
        evidence = ResearchResults(
            hypotheses_evidence=[
                HypothesisEvidence(
                    hypothesis_id=1,
                    evidence=[
                        SearchResult(source="IDC", url="https://idc.com", snippet="Phase 1 launched in 2022.", relevance_score=0.9),
                        SearchResult(source="Gartner", url="https://gartner.com", snippet="Phase 2 grew 30%.", relevance_score=0.8),
                    ],
                    supports=True,
                    confidence="high",
                    conclusion="Supported",
                )
            ],
            total_sources=2,
        )
        path = await gen.create_presentation("Roadmap", storyline, evidence, "short")
        try:
            prs = Presentation(path)
            assert len(prs.slides) >= 4
        finally:
            os.remove(path)


# ---------------------------------------------------------------------------
# Markdown stripping
# ---------------------------------------------------------------------------

class TestMarkdownStripping:

    def test_strip_markdown_removes_bold(self):
        gen = SlideGenerator()
        assert gen._strip_markdown("**Market** grows 20%") == "Market grows 20%"

    def test_strip_markdown_removes_italic(self):
        gen = SlideGenerator()
        assert gen._strip_markdown("*important* finding") == "important finding"

    def test_strip_markdown_removes_links(self):
        gen = SlideGenerator()
        assert gen._strip_markdown("[Gartner](https://gartner.com)") == "Gartner"

    def test_strip_markdown_no_op_on_plain_text(self):
        gen = SlideGenerator()
        assert gen._strip_markdown("Plain text unchanged") == "Plain text unchanged"

    def test_insight_sidebar_strips_markdown(self):
        """_add_insight_sidebar strips **bold** from bullet text."""
        from pptx import Presentation as Prs
        gen = SlideGenerator()
        prs = Prs()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        gen._add_insight_sidebar(
            slide,
            headline="Test",
            bullets=["**Revenue** grew **35%** YoY", "Plain bullet"],
        )
        # Collect all text from the slide
        all_text = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
        assert "**" not in all_text

    async def test_recommendations_strip_markdown(self, sample_storyline, sample_research_results):
        """Recommendation items have ** stripped in the output PPTX."""
        gen = SlideGenerator()
        # Inject bold markers into recs
        sample_storyline.recommendation_items = [
            "**Launch** hybrid cloud pilot by **Q2 2025**",
            "Achieve **SOC2** certification",
        ]
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "short"
        )
        try:
            prs = Presentation(path)
            # Recs slide is second-to-last
            recs_slide = prs.slides[-2]
            texts = " ".join(s.text_frame.text for s in recs_slide.shapes if s.has_text_frame)
            assert "**" not in texts
        finally:
            os.remove(path)


# ---------------------------------------------------------------------------
# Dynamic chart slides
# ---------------------------------------------------------------------------

class TestDynamicChartSlides:

    async def test_medium_only_adds_slides_for_available_slide_data(self):
        """Medium deck with only bar_chart in slide_data adds just 1 extra chart slide."""
        gen = SlideGenerator()
        storyline = _single_hyp_storyline("bar")
        storyline.slide_data = {
            "bar_chart": {
                "action_title": "Cloud leads at 72%",
                "categories": ["Cloud", "On-prem"],
                "values": [72, 28],
                "metric": "Share (%)"
            }
            # No waterfall, pie, tornado
        }
        path = await gen.create_presentation("Cloud", storyline, _empty_research(), "medium")
        try:
            prs = Presentation(path)
            # title(1) + situation(1) + complication(1) + 1 hypothesis + 1 bar chart + recs(1) + sources(1) = 7
            assert len(prs.slides) == 7
        finally:
            os.remove(path)

    async def test_medium_all_four_chart_types_produces_correct_count(
        self, sample_storyline, sample_research_results
    ):
        """Medium deck with all 4 chart types adds 4 extra chart slides (default fixture)."""
        gen = SlideGenerator()
        path = await gen.create_presentation(
            "Cloud Strategy", sample_storyline, sample_research_results, "medium"
        )
        try:
            prs = Presentation(path)
            # title + sit + comp + 3 hyp + bar + waterfall + pie + tornado + recs + sources = 12
            assert len(prs.slides) == 12
        finally:
            os.remove(path)
